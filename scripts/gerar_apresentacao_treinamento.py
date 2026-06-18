# -*- coding: utf-8 -*-
"""
Gera a apresentacao de treinamento do CTG.Engenharia NO PADRAO VISUAL CTG BRASIL.

- Usa o template corporativo docs/_template_ctg.pptx (layout 'BG Claro' com a
  onda colorida no topo, fonte CTG Sans, titulos em azul italico).
- Capa reconstruida com foto da usina + ondas azuis + logo (assets em
  docs/template_assets/).
- Screenshots reais das telas em docs/screenshots/.
- Modulo financeiro (Forecast) fora de escopo.

Uso:   python scripts/gerar_apresentacao_treinamento.py
Saida: docs/Treinamento_CTG_Engenharia.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- paths
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE = os.path.join(ROOT, "docs", "_template_ctg.pptx")
ASSETS   = os.path.join(ROOT, "docs", "template_assets")
IMG      = os.path.join(ROOT, "docs", "screenshots")
OUT      = os.path.join(ROOT, "docs", "Treinamento_CTG_Engenharia.pptx")

# ---------------------------------------------------------------- estilo CTG
FONT        = "CTG Sans"
TITLE_BLUE  = RGBColor(0x00, 0x70, 0xC0)
TEXT        = RGBColor(0x3F, 0x3F, 0x3F)
TEXT_SOFT   = RGBColor(0x5A, 0x63, 0x6E)
BRANCO      = RGBColor(0xFF, 0xFF, 0xFF)
VERDE       = RGBColor(0x15, 0x80, 0x3D)
VERM        = RGBColor(0x99, 0x33, 0x33)
CINZA_CLARO = RGBColor(0xF2, 0xF5, 0xF8)
VERDE_CLARO = RGBColor(0xE7, 0xF4, 0xEC)
AZUL_CLARO  = RGBColor(0xE7, 0xF1, 0xF8)
BORDA       = RGBColor(0xCB, 0xD5, 0xE1)

SLIDE_W = Inches(13.333)

prs = Presentation(TEMPLATE)
# layout corporativo com a onda colorida no topo
BG = [l for l in prs.slide_layouts if l.name == "BG Claro"][0]

# remove os slides de exemplo do template (mantem master/layouts/tema/midia).
# remove tambem a relacao -> a parte fica orfa e nao e gravada (evita partnames duplicados).
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst):
    rId = sldId.get(qn("r:id"))
    if rId:
        prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


# ---------------------------------------------------------------- helpers
def add():
    return prs.slides.add_slide(BG)


def rect(slide, x, y, w, h, color, line=None, line_w=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def tbox(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.06):
    """paras: lista de paragrafos; cada paragrafo = lista de (texto,size,bold,italic,color,level)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(space_after)
        p.space_before = Pt(0); p.line_spacing = line_spacing
        if isinstance(para, tuple):
            para = [para]
        p.level = para[0][5] if len(para[0]) > 5 else 0
        for seg in para:
            text, size, bold, ital, color = seg[0], seg[1], seg[2], seg[3], seg[4]
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size); r.font.bold = bold; r.font.italic = ital
            r.font.color.rgb = color; r.font.name = FONT
    return tb


def title(slide, text):
    """Titulo no padrao CTG: azul, italico, MAIUSCULAS, topo-esquerda (sobre a onda)."""
    tbox(slide, Inches(0.4), Inches(0.16), Inches(12.5), Inches(0.7),
         [[(text.upper(), 23, True, True, TITLE_BLUE)]])


def bullets(slide, items, x, y, w, h=Inches(4.6), size=15, soft_sub=True):
    paras = []
    for it in items:
        text, lvl = it if isinstance(it, tuple) else (it, 0)
        if lvl == 0:
            prefix = "•  "; color = TEXT; sz = size
            bold = text.endswith(":")
        else:
            prefix = "–  "; color = TEXT_SOFT if soft_sub else TEXT; sz = size - 1
            bold = False
        paras.append([(prefix + text, sz, bold, False, color, lvl)])
    tbox(slide, x, y, w, h, paras, space_after=7, line_spacing=1.1)


def place_image(slide, fname, x, y, w):
    h = Emu(int(w * 9 / 16))
    rect(slide, Emu(x - Inches(0.035)), Emu(y - Inches(0.035)),
         Emu(w + Inches(0.07)), Emu(h + Inches(0.07)), BORDA)
    path = os.path.join(IMG, fname)
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, width=w, height=h)
    else:
        rect(slide, x, y, w, h, CINZA_CLARO)
        tbox(slide, x, y, w, h, [[("(%s)" % fname, 12, False, False, TEXT)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return x, y, w, h


def tip_bar(slide, text, y=Inches(6.55), x=Inches(0.4), w=Inches(12.5)):
    rect(slide, x, y, w, Inches(0.6), AZUL_CLARO)
    rect(slide, x, y, Inches(0.07), Inches(0.6), TITLE_BLUE)
    tbox(slide, Emu(x + Inches(0.22)), y, Emu(w - Inches(0.4)), Inches(0.6),
         [[("Dica:  ", 13, True, False, TITLE_BLUE),
           (text, 13, False, False, TEXT)]], anchor=MSO_ANCHOR.MIDDLE)


def module_slide(title_txt, intro, items, image, tip=None):
    s = add()
    title(s, title_txt)
    tbox(s, Inches(0.4), Inches(1.32), Inches(12.5), Inches(0.7),
         [[(intro, 13.5, False, False, TEXT)]])
    bullets(s, items, x=Inches(0.4), y=Inches(2.15), w=Inches(4.55),
            h=Inches(4.2), size=14)
    place_image(s, image, Inches(5.3), Inches(2.0), Inches(7.55))
    if tip:
        tip_bar(s, tip)
    return s


def card(slide, x, y, w, h, nome, itens, cor):
    """Cartao de ferramenta: faixa colorida com nome + lista do que controla."""
    rect(slide, x, y, w, h, BRANCO, line=BORDA)
    rect(slide, x, y, w, Inches(0.5), cor)
    tbox(slide, x, y, w, Inches(0.5),
         [[(nome, 14, True, False, BRANCO)]], align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
    tbox(slide, Emu(x + Inches(0.15)), Emu(y + Inches(0.6)),
         Emu(w - Inches(0.3)), Emu(h - Inches(0.7)),
         [[(itens, 13, False, False, TEXT)]], anchor=MSO_ANCHOR.MIDDLE,
         align=PP_ALIGN.CENTER)


def chip(slide, x, y, w, texto, cor_fundo, cor_txt, dashed=False):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.42))
    sp.fill.solid(); sp.fill.fore_color.rgb = cor_fundo
    if dashed:
        sp.line.color.rgb = cor_txt; sp.line.width = Pt(1)
        sp.line.dash_style = 2  # dash
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = texto
    r.font.size = Pt(12); r.font.bold = True; r.font.name = FONT
    r.font.color.rgb = cor_txt
    return sp


# cores das ferramentas atuais
LISTS_ROXO  = RGBColor(0x6C, 0x2D, 0x7C)
EXCEL_VERDE = RGBColor(0x1D, 0x6F, 0x42)
WORD_AZUL   = RGBColor(0x2B, 0x57, 0x9A)


# ============================================================ 1. CAPA
s = add()
s.shapes.add_picture(os.path.join(ASSETS, "cover_dam.jpg"), 0, 0,
                     width=Inches(13.333), height=Inches(7.48))
s.shapes.add_picture(os.path.join(ASSETS, "cover_filter.png"), 0, 0,
                     width=Inches(10.95), height=Inches(7.5))
s.shapes.add_picture(os.path.join(ASSETS, "cover_waves.png"), 0, Inches(3.67),
                     width=Inches(13.333), height=Inches(3.95))
s.shapes.add_picture(os.path.join(ASSETS, "logo.png"), Inches(0.48), Inches(0.45),
                     width=Inches(3.13), height=Inches(0.86))
_NAVY = RGBColor(0x0B, 0x3A, 0x5E)
tbox(s, Inches(0.55), Inches(2.7), Inches(5.6), Inches(2.4), [
    [("CTG.Engenharia", 32, True, False, _NAVY)],
    [("Treinamento da plataforma", 18, False, False, _NAVY)],
    [("Time de Engenharia", 18, False, False, _NAVY)],
    [("Investimentos · contratos · equipamentos · documentos · pessoas",
      12, False, True, RGBColor(0x2A, 0x5A, 0x7E))],
], space_after=5)
tbox(s, Inches(8.5), Inches(6.95), Inches(4.4), Inches(0.4),
     [[("Junho, 2026  ·  Forecast fora de escopo", 12, True, False, BRANCO)]],
     align=PP_ALIGN.RIGHT)

# ============================================================ 2. AGENDA
s = add()
title(s, "O que vamos ver hoje")
bullets(s, [
    "O que e o CTG.Engenharia",
    "Por que migrar: hoje x plataforma",
    "Acesso, perfis e navegacao",
    "IACs 2026 — autorizacoes de investimento",
    "Acompanhamento de Projetos (contratos)",
    "Mapa de Equipamentos",
], x=Inches(0.6), y=Inches(1.6), w=Inches(6.0), size=17)
bullets(s, [
    "Controle de Documentos",
    "Cronograma (Project) e Metas",
    "Ferias e Delegacao de acesso",
    "Relatorios e exportacoes",
    "Instalacao no celular (PWA)",
], x=Inches(6.9), y=Inches(1.6), w=Inches(6.0), size=17)

# ============================================================ 3. O QUE E
s = add()
title(s, "O que e o CTG.Engenharia")
tbox(s, Inches(0.4), Inches(1.3), Inches(12.4), Inches(0.9),
     [[("Plataforma web da CTG Brasil que centraliza a gestao das atividades de engenharia de manutencao — substituindo planilhas soltas por dados unificados, com graficos, historico e acesso pelo celular.", 16, False, False, TEXT)]])
bullets(s, [
    "Investimentos (IACs) — autorizacoes de investimento",
    "Contratos e projetos em andamento — valores, status, cronograma",
    "Equipamentos das subestacoes — inventario tecnico",
    "Documentos tecnicos — codigo padronizado e revisoes",
    "Cronogramas, Metas e Ferias da equipe",
    "Web + aplicativo no celular (PWA) — funciona em qualquer tela",
], x=Inches(0.5), y=Inches(2.5), w=Inches(12), size=16)

# ============================================================ 4. COMO E FEITO HOJE
s = add()
title(s, "Como e feito hoje — ferramentas dispersas")
tbox(s, Inches(0.4), Inches(1.3), Inches(12.4), Inches(0.6),
     [[("Hoje o controle esta espalhado em tres ferramentas diferentes, sem conversa entre si:", 15, False, False, TEXT)]])
card(s, Inches(0.5),  Inches(2.0), Inches(3.9), Inches(1.6),
     "Microsoft Lists", "IACs  ·  Acomp. de Projetos", LISTS_ROXO)
card(s, Inches(4.72), Inches(2.0), Inches(3.9), Inches(1.6),
     "Excel", "Ferias  ·  Metas", EXCEL_VERDE)
card(s, Inches(8.94), Inches(2.0), Inches(3.9), Inches(1.6),
     "Word", "Documentos (numeracao manual)", WORD_AZUL)
# barra unificada
rect(s, Inches(0.5), Inches(3.95), Inches(12.34), Inches(0.62), TITLE_BLUE)
tbox(s, Inches(0.5), Inches(3.95), Inches(12.34), Inches(0.62),
     [[("↓   CTG.Engenharia — uma plataforma, um acesso, uma fonte de verdade", 15, True, False, BRANCO)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, [
    "Silos de dados — sem visao consolidada",
    "Sem padrao de acesso — cada arquivo com seu compartilhamento",
    "Retrabalho e versoes duplicadas por e-mail",
    "Numeracao e padronizacao manuais (ex.: codigo no Word)",
    "Sem historico unificado nem indicadores automaticos",
], x=Inches(0.5), y=Inches(4.85), w=Inches(12.3), size=14)

# ============================================================ 4b. VANTAGENS DA MIGRACAO
s = add()
title(s, "Vantagens da migracao")
rect(s, Inches(0.6), Inches(1.55), Inches(5.85), Inches(5.0), CINZA_CLARO)
rect(s, Inches(6.9), Inches(1.55), Inches(5.85), Inches(5.0), VERDE_CLARO)
tbox(s, Inches(0.6), Inches(1.68), Inches(5.85), Inches(0.5),
     [[("HOJE — Lists + Excel + Word", 15, True, False, VERM)]], align=PP_ALIGN.CENTER)
tbox(s, Inches(6.9), Inches(1.68), Inches(5.85), Inches(0.5),
     [[("COM O CTG.ENGENHARIA", 15, True, False, VERDE)]], align=PP_ALIGN.CENTER)
bullets(s, [
    "Dados em 3 ferramentas separadas",
    "Compartilhamento por arquivo",
    "Sem visao consolidada",
    "Numeracao e padroes manuais",
    "Sem rastreabilidade",
    "Informacao isolada",
    "Acesso so no computador",
], x=Inches(0.85), y=Inches(2.35), w=Inches(5.4), size=14)
bullets(s, [
    "Tudo em um so lugar, login unico",
    "Controle de acesso por perfil",
    "KPIs e graficos automaticos",
    "Codigos automaticos e validacao",
    "Historico, responsaveis e check-ins",
    "Dados relacionados entre modulos",
    "Web + app no celular (PWA)",
], x=Inches(7.15), y=Inches(2.35), w=Inches(5.4), size=14)

# ============================================================ 4c. PRONTO PARA CRESCER
s = add()
title(s, "Pronto para crescer — integracao de novas ferramentas")
tbox(s, Inches(0.4), Inches(1.3), Inches(12.4), Inches(0.9),
     [[("A maior vantagem nao e so substituir o que existe, mas abrir caminho para o que vem: a plataforma e unica e extensivel — novos modulos entram no mesmo ambiente.", 15, False, False, TEXT)]])
bullets(s, [
    "Novos modulos e integracoes entram no mesmo lugar",
    "Reaproveitam login, perfis, usinas e usuarios",
    "Mais controle dos processos ponta a ponta",
    "Evolucao governada e auditavel — nao 'mais um Excel'",
], x=Inches(0.5), y=Inches(2.35), w=Inches(6.2), size=15)
# chips: modulos existentes + futuros
tbox(s, Inches(7.0), Inches(2.2), Inches(5.8), Inches(0.4),
     [[("JA NA PLATAFORMA", 12, True, False, TITLE_BLUE)]])
exist = ["IACs", "Projetos", "Equipamentos", "Documentos", "Cronograma", "Metas", "Ferias"]
cx, cy = Inches(7.0), Inches(2.65)
for i, m in enumerate(exist):
    col = i % 3; row = i // 3
    chip(s, Emu(cx + col * Inches(1.95)), Emu(cy + row * Inches(0.55)),
         Inches(1.8), m, AZUL_CLARO, TITLE_BLUE)
tbox(s, Inches(7.0), Inches(4.55), Inches(5.8), Inches(0.4),
     [[("EVOLUCAO FUTURA", 12, True, False, VERDE)]])
fut = ["+ Nova ferramenta", "+ Integracao", "+ Novo processo"]
fy = Inches(5.0)
for i, m in enumerate(fut):
    col = i % 2; row = i // 2
    chip(s, Emu(cx + col * Inches(2.95)), Emu(fy + row * Inches(0.55)),
         Inches(2.8), m, BRANCO, VERDE, dashed=True)
tip_bar(s, "Quanto mais centralizado, maior o controle e a visao do conjunto.")

# ============================================================ 4d. PONTOS DE ATENCAO
s = add()
title(s, "Pontos de atencao (sejamos honestos)")
tbox(s, Inches(0.4), Inches(1.3), Inches(12.4), Inches(0.5),
     [[("Toda migracao tem trade-offs — conhece-los ajuda a planejar a transicao:", 15, False, False, TEXT)]])
# cabecalho de duas colunas
rect(s, Inches(0.5), Inches(1.95), Inches(6.0), Inches(0.5), VERM)
rect(s, Inches(6.7), Inches(1.95), Inches(6.1), Inches(0.5), VERDE)
tbox(s, Inches(0.5), Inches(1.95), Inches(6.0), Inches(0.5),
     [[("PONTO DE ATENCAO", 13, True, False, BRANCO)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
tbox(s, Inches(6.7), Inches(1.95), Inches(6.1), Inches(0.5),
     [[("COMO MITIGAR", 13, True, False, BRANCO)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
pares = [
    ("Curva de aprendizado / mudanca", "Treinamento, Tutorial e Manual"),
    ("Migracao dos dados historicos", "Importadores com preview ja prontos"),
    ("Disciplina para manter atualizado", "Alertas de desatualizacao e check-ins"),
    ("Depende de conexao (web)", "PWA no celular, acesso de qualquer lugar"),
    ("Analises ad-hoc no Excel", "Exportacao para Excel em todas as telas"),
    ("Governanca (acessos e backup)", "Administracao central de perfis e dados"),
]
ry = Inches(2.55)
for i, (a, b) in enumerate(pares):
    bg = BRANCO if i % 2 == 0 else CINZA_CLARO
    rect(s, Inches(0.5), ry, Inches(6.0), Inches(0.62), bg, line=BORDA)
    rect(s, Inches(6.7), ry, Inches(6.1), Inches(0.62), bg, line=BORDA)
    tbox(s, Emu(Inches(0.5) + Inches(0.15)), ry, Inches(5.7), Inches(0.62),
         [[(a, 13, False, False, TEXT)]], anchor=MSO_ANCHOR.MIDDLE)
    tbox(s, Emu(Inches(6.7) + Inches(0.15)), ry, Inches(5.8), Inches(0.62),
         [[(b, 13, False, False, TEXT)]], anchor=MSO_ANCHOR.MIDDLE)
    ry = Emu(ry + Inches(0.66))

# ============================================================ 5. ACESSO + LOGIN
s = add()
title(s, "Acesso e perfis")
bullets(s, [
    "Acesso: e-mail institucional + senha",
    "Esqueceu? 'Esqueceu sua senha?' envia o link",
    "Tambem da para entrar com Microsoft",
    "Perfis: Administrador, Planejador,",
    ("Coordenador, Gestor e Engenheiro", 1),
    "O menu muda conforme o perfil",
], x=Inches(0.4), y=Inches(1.6), w=Inches(4.55), size=15)
place_image(s, "login.png", Inches(5.3), Inches(2.0), Inches(7.55))
tip_bar(s, "Cada perfil ve menus diferentes — por isso alguns itens aparecem so para alguns usuarios.")

# ============================================================ 6. NAVEGACAO
s = add()
title(s, "O menu lateral por secoes")
nav = (
    "Inicio  —  visao geral\n\n"
    "LISTS\n   Acomp. de Projetos\n   IACs 2026\n\n"
    "MAPA DE EQUIPAMENTOS\n   Mapa de Equipamentos\n   Gestao de Equipamentos\n\n"
    "GESTAO DE PESSOAS E DOCUMENTOS\n   Ferias · Documentos\n   Metas · Cronograma Project\n\n"
    "ADMINISTRACAO\n   Gerenciar Usuarios\n\n"
    "Rodape: Tutorial · Sugestoes ·\nConfiguracoes · Meu Perfil · Sair"
)
rect(s, Inches(0.6), Inches(1.55), Inches(6.2), Inches(5.5), CINZA_CLARO)
tb = s.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(5.7), Inches(5.2))
tf = tb.text_frame; tf.word_wrap = True; tf.paragraphs[0].text = nav
for p in tf.paragraphs:
    for r in p.runs:
        r.font.name = "Consolas"; r.font.size = Pt(13); r.font.color.rgb = TEXT
bullets(s, [
    "No celular: menu inferior + botao de menu",
    "O conteudo depende do seu perfil",
    "Use os filtros por coluna nas tabelas",
    "Sino de notificacoes no topo (alertas e mensagens)",
], x=Inches(7.2), y=Inches(2.1), w=Inches(5.5), size=16)

# ============================================================ 7. IACs
module_slide(
    "IACs 2026 — Autorizacoes de Investimento",
    "Registram solicitacoes de investimento em equipamentos, servicos e obras, e acompanham o ciclo ate a contratacao.",
    [
        "3 partes: graficos de Status e",
        ("Prioridade + lista filtravel", 1),
        "Clique nas barras/rosca para filtrar",
        "Ciclo: Not started -> ... -> Contract signed",
        "Abas: Todos · Meus IACs · por area",
        "Novo, Editar, Check-in",
        "Importar/Exportar Excel",
    ],
    "iacs.png",
    tip="Importar Excel mostra um preview (novos / atualizados / pulados) antes de confirmar.",
)

# ============================================================ 8. ACOMPANHAMENTO
module_slide(
    "Acompanhamento de Projetos",
    "Acompanhamento de contratos e projetos em andamento: valores, saldo, status e cronograma.",
    [
        "Graficos: Natureza (CAPEX/OPEX/",
        ("Guarda-chuva) e por Usina", 1),
        "Status: andamento, encerramento,",
        ("paralisado, capitalizado...", 1),
        "Valores: Contrato, Realizado, Saldo + SI",
        ("Saldo automatico; negativo em vermelho", 1),
        "Gerar Relatorio HTML mensal",
    ],
    "acompanhamento.png",
    tip="O Relatorio HTML mensal pode ser gerado a partir do banco ou de um Excel enviado.",
)

# ============================================================ 9. EQUIPAMENTOS
module_slide(
    "Mapa e Gestao de Equipamentos",
    "Inventario tecnico dos equipamentos das subestacoes, navegavel por usina, funcao e tipo.",
    [
        "Mapa de Equipamentos (consulta):",
        ("Selecione Usina -> Funcao -> Tipo", 1),
        ("Veja distribuicao e lista detalhada", 1),
        "Gestao de Equipamentos (admin):",
        ("Cadastrar, editar e manter os dados", 1),
    ],
    "equipamentos.png",
    tip="O Mapa e para consultar; a Gestao e onde os dados sao mantidos atualizados.",
)

# ============================================================ 10. DOCUMENTOS
module_slide(
    "Controle de Documentos",
    "Documentos tecnicos com codigo padronizado, controle de revisoes e status.",
    [
        "Tipos: ATA, RT, ET, MC, ROP...",
        "Codigo automatico:",
        ("TIPO-AREA-SEQ-ANO[-RREV]", 1),
        "Status: elaboracao, aprovacao,",
        ("publicado (link), cancelado", 1),
        "Nova Revisao mantem o original",
        "Importar .docx; Exportar Excel/HTML",
    ],
    "documentos.png",
    tip="Publicado exige o link do documento; Cancelado e irreversivel sem permissao.",
)

# ============================================================ 11. CRONOGRAMA
module_slide(
    "Cronograma (Project)",
    "Cronograma executivo estilo Gantt para atividades de engenharia e de campo.",
    [
        "Fases e tarefas por WBS (1, 1.1...)",
        "Datas de inicio/fim e progresso (%)",
        "Dependencias entre tarefas",
        "Visao por Dia / Semana / Mes",
        "Caminho critico e impressao em PDF",
        "Ideal p/ comissionamentos e campo",
    ],
    "cronograma.png",
    tip="Use + Ativ. Macro e + Atividade para montar o cronograma; vincule por dependencias.",
)

# ============================================================ 12. METAS
module_slide(
    "Metas",
    "Acompanhamento das metas da equipe, com pesos e percentual de atingimento por colaborador.",
    [
        "Metas coletivas e individuais",
        "Cada meta tem um peso (importancia)",
        "Valor atingido -> % de atingimento",
        "Atingimento ponderado por pessoa",
        "Relatorio individual por colaborador",
        "KPIs de metas cadastradas/concluidas",
    ],
    "metas.png",
    tip="O peso define a importancia relativa de cada meta no resultado final.",
)

# ============================================================ 13. FERIAS
module_slide(
    "Controle de Ferias",
    "Planejamento das ferias da equipe com timeline visual e alerta de conflitos na mesma area.",
    [
        "Escolha Ano e Area",
        "KPIs: colaboradores, marcados,",
        ("periodos no ADP, total de dias", 1),
        "Timeline: 1 linha por pessoa;",
        ("linha vermelha = hoje", 1),
        "Ate 3 periodos por pessoa",
        "Adicionar / Editar / Excluir",
    ],
    "ferias.png",
    tip="O sistema alerta quando duas pessoas da mesma area tiram ferias no mesmo periodo.",
)

# ============================================================ 14. DELEGACAO
s = add()
title(s, "Delegacao de Acesso")
tbox(s, Inches(0.4), Inches(1.3), Inches(12.4), Inches(0.8),
     [[("Permite que outra pessoa assuma seus projetos temporariamente — por exemplo, durante suas ferias.", 16, False, False, TEXT)]])
bullets(s, [
    "Meu Perfil -> Delegacao de Acesso -> + Nova Delegacao",
    "Selecione o usuario, defina inicio e fim e (opcional) o motivo",
    "Ativa e desativa sozinha nas datas definidas",
    "Pode ser revogada a qualquer momento",
    "Acoes ficam registradas no nome do delegado, nao no seu",
    "Coordenadores, planejadores e engenheiros podem delegar",
], x=Inches(0.5), y=Inches(2.4), w=Inches(12), size=17)

# ============================================================ 15. RELATORIOS
s = add()
title(s, "Relatorios e Exportacoes")
bullets(s, [
    "Relatorio HTML: escolha Mes e Ano e a fonte (banco ou Excel) -> Gerar",
    "Exportacao Excel em: IACs, Acompanhamento, Ferias, Documentos e outras",
    "Documentos tambem exporta em HTML com estatisticas",
    "Importacoes sempre tem tela de preview com estatisticas antes de confirmar",
], x=Inches(0.5), y=Inches(1.6), w=Inches(12), size=17)
tip_bar(s, "Confira sempre o preview de importacao antes de confirmar os dados.", y=Inches(4.8))

# ============================================================ 16. ADMIN / USUARIOS
module_slide(
    "Gestao de Usuarios",
    "Administradores criam e gerenciam os usuarios, perfis e areas — e aprovam solicitacoes de acesso.",
    [
        "Criar / editar / inativar usuarios",
        "Atribuir perfil e area",
        "Aprovar solicitacoes de acesso",
        "Redefinir senha de usuarios",
        "KPIs por tipo de perfil",
    ],
    "inicio.png",
    tip="Cada usuario recebe um perfil que define exatamente o que pode ver e fazer.",
)

# ============================================================ 17. RECURSOS + PWA
s = add()
title(s, "Recursos gerais e app no celular")
bullets(s, [
    "Notificacoes (sino): alertas e mensagens nao lidas",
    "Tutorial interativo e canal de Sugestoes/Feedback",
    "Check-ins (marque o que ja revisou) e filtros por coluna",
    "Admin: criar/editar usuarios, perfis, areas e acessos",
], x=Inches(0.5), y=Inches(1.6), w=Inches(7.4), size=16)
rect(s, Inches(8.3), Inches(1.7), Inches(4.4), Inches(4.5), VERDE_CLARO)
tbox(s, Inches(8.55), Inches(1.85), Inches(3.9), Inches(0.5),
     [[("Instalar no celular (PWA)", 15, True, False, VERDE)]])
bullets(s, [
    "iPhone (Safari): Compartilhar -> Adicionar a Tela de Inicio",
    "Android (Chrome): tres pontos -> Instalar aplicativo",
    "Vira app em tela cheia, com badge e aviso de nova versao",
], x=Inches(8.55), y=Inches(2.5), w=Inches(3.9), size=13)

# ============================================================ 18. EXERCICIOS
s = add()
title(s, "Exercicios sugeridos")
bullets(s, [
    "IAC: buscar um IAC, filtrar por area e fazer um check-in",
    "Acompanhamento: abrir um contrato, conferir o saldo e gerar relatorio HTML",
    "Documento: criar um novo e observar o codigo automatico; criar uma revisao",
    "Ferias: registrar um periodo e localizar um conflito na timeline",
    "Delegacao: criar uma delegacao com inicio/fim e depois revogar",
    "Mobile: instalar o app no celular (PWA)",
], x=Inches(0.5), y=Inches(1.6), w=Inches(12), size=17)

# ============================================================ 19. ENCERRAMENTO
s = add()
s.shapes.add_picture(os.path.join(ASSETS, "cover_waves.png"), 0, Inches(3.67),
                     width=Inches(13.333), height=Inches(3.95))
s.shapes.add_picture(os.path.join(ASSETS, "logo.png"), Inches(0.5), Inches(0.6),
                     width=Inches(3.13), height=Inches(0.86))
tbox(s, Inches(0.6), Inches(2.0), Inches(12), Inches(1.0),
     [[("Obrigado!", 40, True, False, TITLE_BLUE)]])
tbox(s, Inches(0.62), Inches(3.0), Inches(11.5), Inches(1.2),
     [[("Duvidas? Use o menu Sugestoes dentro do sistema ou fale com o administrador.", 16, False, False, TEXT)],
      [("Consulte tambem o Tutorial interativo e o Manual do Usuario.", 14, False, True, TEXT_SOFT)]])

# ---------------------------------------------------------------- salvar
prs.save(OUT)
print("OK ->", OUT)
print("Slides:", len(prs.slides._sldIdLst))
